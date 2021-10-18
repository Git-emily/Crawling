from datetime import time

import pythoncom
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT,MSO_VERTICAL_ANCHOR

from pptx import Presentation
from pptx.util import Pt,Cm,Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION,XL_TICK_MARK,XL_DATA_LABEL_POSITION
from win32com.client import DispatchEx


def export_PDF(inputdata):
    DATE = time.strftime("%Y%m%d%H%M",time.localtime(inputdata['Time']))
    JOB_SITE = inputdata['Location_Info']
    ppt = DispatchEx('PowerPoint.Application')  #Start a new Thread
    # ppt.Visible = 0  #back_end running
    ppt.DisplayAlerts = 0 # No Alarm/Alert
    pptSel = ppt.Presentations.Open(app_path + '\\static\\exportedReports\\'+'HBpulse.pptx',WithWindow = False) #back_end running
    # pptSel.Slides(1).Copy()
    # pptSel.Slides.Paste()
    # pptSel.SaveAs('D:\\Apache24\\branch_V3\\App\\static\\exportedReports\\' + 'Building_' + 'HBpulse.pptx')
    pptSel.SaveAs(app_path + '\\static\\exportedReports\\'+DATE+JOB_SITE+'_HBpulse.pptx')
    pptSel.Close()
    add_Content(DATE,JOB_SITE,inputdata)
    ppt.Quit()

def add_Content(DATE,JOB_SITE,inputdata):
    ppt = Presentation(app_path + '\\static\\exportedReports\\'+DATE+JOB_SITE+'_HBpulse.pptx')
    #PPT Slide 1
    slide = ppt.slides[0]
    left,top,width,height = Cm(1.27),Cm(14.74),Cm(10),Cm(2)
    paragraph = slide.shapes.add_textbox(left,top,width,height).text_frame
    multiLineContent ='COUNTRY:  '+ inputdata['Country']+'\n'+'JOB SITE:  '+ inputdata['Location_Info']+'\n'+'ZONE:  '+ inputdata['Zone_Info']+'\n'
    paragraph.paragraphs[0].text = multiLineContent
    paragraph.paragraphs[0].font.size = Pt(18)
    paragraph.paragraphs[0].font.name = 'Arial'
    paragraph.paragraphs[0].font.color.rgb = RGBColor(21,44,115)
    paragraph.paragraphs[0].font.bold = True

    # PPT_Silde'Survey_Info'
    slide = ppt.slides[2]
    left,top,width,height = Cm(1),Cm(4),Cm(13.6),Cm(8)
    shape = slide.shapes.add_table(10,3,left,top,width,height)
    table = shape.table
    table.columns[0].width = Cm(4)
    table.columns[1].width = Cm(6)
    table.columns[2].width = Cm(3.6)

    table.rows[0].height = Cm(2)
    try:
        table.cell(1,0).merge(table.cell(3,0))
        table.cell(4,0).merge(table.cell(6,0))
        table.cell(7,0).merge(table.cell(9,0))
    except Exception as e:
        print('E:',e)
    table.cell(0,0).text = 'Principle'
    table.cell(0,1).text = 'Items'
    table.cell(0,2).text = 'Results'
    table.cell(1, 0).text = 'Building Info'
    table.cell(4, 0).text = 'Ventilation'
    table.cell(7, 0).text = 'Filters'

    content_arry = [["Area*Height",str(float(inputdata['Area'])*float(inputdata['Height']))],["Occupants",inputdata['Occupants']],
                    ["OA CO2",inputdata['OA_CO2']],["OA Flow Rate",inputdata['OA_Flow_Rate']],["OA ACH",inputdata['OA_ACH_Input']],
                    ["CFM/Persom",inputdata['CFM_Person']],["AHU Filter",inputdata['AHU_Filter']],["Air Purifi",inputdata['select_AirPur']],
                    ["UVGI",inputdata['select_UVGI_Input']]]
    for rows in range(9):
        for cols in range(2):
            table.cell(rows+1,cols+1).text = content_arry[rows][cols]

    # PPT_Silde 3,5 add table
    n_Page = [3,5]
    for i in n_Page:
        slide = ppt.slides[i]
        left,top,width,height = Cm(1),Cm(4),Cm(12.6),Cm(8)
        shape = slide.shapes.add_table(12,3,left,top,width,height)
        table = shape.table
        table.columns[0].width = Cm(4)
        table.columns[1].width = Cm(5)
        table.columns[2].width = Cm(3.6)
        table.rows[0].height = Cm(2)

        table.cell(1,0).merge(table.cell(4,0))
        table.cell(5,0).merge(table.cell(8,0))
        table.cell(9,0).merge(table.cell(11,0))

        table.cell(0,0).text = 'Principle'
        table.cell(0,1).text = 'Items'
        table.cell(0,2).text = 'Results'
        table.cell(1, 0).text = 'Gases'
        table.cell(5, 0).text = 'Particles'
        table.cell(9, 0).text = 'Comfort'

        if i == 3:
            content_arry1 = [["CO2", inputdata['CO2']],["CO", inputdata['CO']],["HCHO", inputdata['HCHO']],
                            ["TVOC", inputdata['TVOC']],["PM2.5", inputdata['PM25']],
                            ["PM10", inputdata['PM10']], ["Fungi", inputdata['Fungi']],
                            ["Bacteria", inputdata['Bacteria']],
                            ["Temperature", inputdata['Temperature']],["Humidity", inputdata['Humidity']],
                            ["OA ACH", inputdata['OA_ACH_Current']]]
            for rows in range(11):
                for cols in range(2):
                    table.cell(rows+1,cols+1).text = content_arry1[rows][cols]

            #bar_chart
            shapes = slide.shapes
            # shapes.title.text = 'Indoor Air Quality Indexes'
            x = ['Pollution', 'Comfort', 'Air Quality', 'Airborne Transmission Risk']
            y = [float(inputdata['Current_IANPI']), float(inputdata['Current_ICI']), float(inputdata['Current_IAQ']),
                 float(inputdata['Current_VTRI'])]

            chart_data = CategoryChartData()
            chart_data.categories = x
            chart_data.add_series('Indoor Air Quality Indexes',values=y)
            left,top,width,height = Cm(15.6),Cm(5),Cm(17),Cm(10)
            graphic_frame = shapes.add_chart(chart_type = XL_CHART_TYPE.BAR_CLUSTERED,x = left,y = top,cx=width,cy=height,chart_data=chart_data)
            chart = graphic_frame.chart
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
            data_labels.font.name = 'Calibri'
            data_labels.font.size = Pt(14)

            # seaborn.set_style('darkgrid',{'font.sans-serif':['Arial','Calibri']})
            # warnings.filterwarnings('ignore')
            # plt.figure(dpi=100)
            # x = ['Pollution','Comfort','Air Quality','Airborne Transmission Risk']
            # y = [float(inputdata['Current_IANPI']),float(inputdata['Current_ICI']),float(inputdata['Current_IAQ']),float(inputdata['Current_VTRI'])]
            # seaborn.barplot(y,x)
            # plt.show()

        else:
            content_arry2 = [["CO2", inputdata['CO2_new']], ["CO", inputdata['CO_new']], ["HCHO", inputdata['HCHO_new']],
                             ["TVOC", inputdata['TVOC_new']], ["PM2.5", inputdata['PM25_new']],
                             ["PM10", inputdata['PM10_new']], ["Fungi", inputdata['Fungi_new']],
                             ["Bacteria", inputdata['Bacteria_new']],
                             ["Temperature", inputdata['Temperature_new']], ["Humidity", inputdata['Humidity_new']],
                             ["OA ACH", inputdata['OA_ACH_new']]]
            for rows in range(11):
                for cols in range(2):
                    table.cell(rows+1,cols+1).text = content_arry2[rows][cols]

            shapes = slide.shapes
            x = ['Pollution', 'Comfort', 'Air Quality', 'Airborne Transmission Risk']
            y = [float(inputdata['Improved_IANPI']), float(inputdata['Improved_ICI']), float(inputdata['Improved_IAQ']),
                 float(inputdata['Improved_VTRI'])]
            chart_data = CategoryChartData()
            chart_data.categories = x
            chart_data.add_series('Indoor Air Quality Indexes_Improved', values=y)
            left, top, width, height = Cm(15.6), Cm(5), Cm(17), Cm(10)
            graphic_frame = shapes.add_chart(chart_type=XL_CHART_TYPE.BAR_CLUSTERED, x=left, y=top, cx=width, cy=height,
                                             chart_data=chart_data)
            chart = graphic_frame.chart
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
            data_labels.font.name = 'Calibri'
            data_labels.font.size = Pt(14)

    #PPT Slide 4 add Table
    slide = ppt.slides[4]
    left,top,width,height = Cm(3),Cm(8),Cm(26),Cm(6)
    shape = slide.shapes.add_table(10,4,left,top,width,height)
    table = shape.table
    table.columns[0].width = Cm(5)
    table.columns[1].width = Cm(6)
    table.columns[2].width = Cm(7)
    table.columns[3].width = Cm(8)
    table.rows[0].height = Cm(1.2)

    table.cell(1,0).merge(table.cell(3,0))
    table.cell(4,0).merge(table.cell(6,0))
    table.cell(7,0).merge(table.cell(9,0))

    table.cell(0,0).text = 'Operations'
    table.cell(0,1).text = 'Priority1(Basic)'
    table.cell(0,2).text = 'Priority2(Advanced)'
    table.cell(0,3).text = 'Status'

    table.cell(1,0).text = 'Ventilation'
    table.cell(4,0).text = 'Filtration'
    table.cell(7,0).text = 'Purification'

    content_arry = [["Outside Air ACH ", '/',inputdata['Impro_ACH']],['/','/','/'],['/','/','/'],
                    ["Air Cleaning Equip", '/',inputdata['select_MERV']],["AHU Filteration", '/',inputdata['select_AHUFilter']],
                    ["UVGI",'/', inputdata['select_UVGI']],['/','/','/'],['/','/','/'],['/','/','/']]

    for rows in range(9):
        for cols in range(3):
            table.cell(rows+1,cols+1).text = content_arry[rows][cols]

    ppt.save(app_path + '\\static\\exportedReports\\'+DATE+JOB_SITE+'_HBpulse.pptx')

if __name__ == '__main__':
    pythoncom.CoInitialize()  # run this before running the new Thread
    export_PDF(inputdata)